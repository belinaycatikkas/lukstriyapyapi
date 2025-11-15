from pptx import Presentation
import os
import sys
import re

# UTF-8 encoding ayarla
if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')

def clean_text(text):
    """Özel karakterleri temizle ve güvenli hale getir"""
    # Zero-width space ve benzeri karakterleri kaldır
    text = re.sub(r'[\u200b-\u200f\u202a-\u202e]', '', text)
    # Diğer kontrol karakterlerini kaldır
    text = re.sub(r'[\x00-\x1f\x7f-\x9f]', '', text)
    return text.strip()

# Dosya yolu
file_path = r"C:\Users\M_ilg\Desktop\belinayodev\urla ppt.pptx"

try:
    # PowerPoint dosyasını aç
    prs = Presentation(file_path)
    
    print("=" * 60)
    print("POWERPOINT SUNUM ANALIZI")
    print("=" * 60)
    print(f"\nDosya: {os.path.basename(file_path)}")
    print(f"Toplam Slayt Sayisi: {len(prs.slides)}")
    print(f"Slayt Boyutu: {prs.slide_width / 914400:.2f} x {prs.slide_height / 914400:.2f} inc")
    print("\n" + "=" * 60)
    
    # Her slaytı incele
    for i, slide in enumerate(prs.slides, 1):
        print(f"\n--- SLAYT {i} ---")
        
        # Slayt düzeni bilgisi
        if slide.slide_layout.name:
            print(f"Duzen: {slide.slide_layout.name}")
        
        # Tüm şekilleri incele
        shapes_with_text = []
        images_count = 0
        
        for shape in slide.shapes:
            # Metin içeren şekiller
            if hasattr(shape, "text") and shape.text.strip():
                cleaned_text = clean_text(shape.text)
                if cleaned_text:
                    shapes_with_text.append(cleaned_text)
            
            # Görsel kontrolü
            if hasattr(shape, "image"):
                images_count += 1
                try:
                    img_filename = shape.image.filename if hasattr(shape.image, 'filename') else 'Bilinmeyen'
                    print(f"  [Gorsel bulundu: {img_filename}]")
                except:
                    print(f"  [Gorsel bulundu: Bilinmeyen format]")
        
        # Metin içeriğini göster
        if shapes_with_text:
            print("Metin Icerigi:")
            for idx, text in enumerate(shapes_with_text, 1):
                # Çok uzun metinleri kısalt
                display_text = text[:200] + "..." if len(text) > 200 else text
                try:
                    print(f"  {idx}. {display_text}")
                except UnicodeEncodeError:
                    # Eğer hala encoding sorunu varsa, ASCII'ye çevir
                    safe_text = display_text.encode('ascii', 'ignore').decode('ascii')
                    print(f"  {idx}. {safe_text}")
        else:
            print("Metin icerigi bulunamadi.")
        
        if images_count > 0:
            print(f"Toplam gorsel sayisi: {images_count}")
        
        # Tablo kontrolü
        tables = [s for s in slide.shapes if hasattr(s, "table")]
        if tables:
            print(f"Tablo sayisi: {len(tables)}")
            for table_idx, table in enumerate(tables, 1):
                print(f"  Tablo {table_idx}: {len(table.table.rows)} satir x {len(table.table.columns)} sutun")
    
    print("\n" + "=" * 60)
    print("ANALIZ TAMAMLANDI")
    print("=" * 60)
    
except ImportError:
    print("HATA: python-pptx kutuphanesi yuklu degil.")
    print("Yuklemek icin: pip install python-pptx")
except Exception as e:
    print(f"HATA: {type(e).__name__}: {str(e)}")
    import traceback
    traceback.print_exc()

